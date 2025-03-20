from crewai.tools import BaseTool
from typing import Type, Literal
from pydantic import BaseModel, Field

class TaskStatusUpdateSchema(BaseModel):
    task_name: str = Field(..., description="The name of the task being updated.")
    status: Literal["COMPLETED", "PENDING"] = Field(..., description="Task status must be 'COMPLETED' or 'PENDING'.")
    chat_message_id: int = Field(..., description="The ID of the chat message associated with the task.")

class TaskStatusUpdate(BaseTool):
    name: str = "task_status_update"
    description: str = "Updates the task status dynamically to either 'COMPLETED' or 'PENDING' along with the task name and chat message ID."
    args_schema: Type[BaseModel] = TaskStatusUpdateSchema  # Ensure correct argument validation

    def _run(self, task_name: str, status: str, chat_message_id: int) -> str:
        """
        Updates the task status of the given task_name associated with the chat_message_id.

        Args:
            task_name (str): The name of the task being updated.
            status (str): Must be 'COMPLETED' or 'PENDING'.
            chat_message_id (int): The ID of the chat message associated with the task.

        Returns:
            str: A message confirming the update.
        """

        # Print and return the result
        print(f"Updated Task: {task_name} | Status: {status} | Chat Message ID: {chat_message_id}")
        return f"Task '{task_name}' status updated successfully to {status} for Chat Message ID {chat_message_id}"


from crewai.tools import BaseTool
from typing import Type
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches


class PresentationDesignerInput(BaseModel):
    """Input schema for PresentationDesigner."""
    product_name: str = Field(..., description="Name of the product being presented.")
    key_features: list[str] = Field(..., description="List of key features of the product.")
    benefits: list[str] = Field(..., description="List of benefits for the customers.")
    testimonials: list[str] = Field(..., description="Customer testimonials to include in the slides.")
    output_file: str = Field(..., description="Output file name for the presentation (e.g., 'presentation.pptx').")
    chat_message_id: int = Field(..., description="The ID of the chat message associated with the task.")


class PresentationDesignerTool(BaseTool):
    name: str = "Presentation Designer"
    description: str = (
        "Generates a 10-slide PowerPoint deck showcasing the product’s key features, benefits, and customer success stories. "
        "Ensures a professional and visually appealing format."
    )
    args_schema: Type[BaseModel] = PresentationDesignerInput

    def _run(self, product_name: str, key_features: list[str], benefits: list[str], testimonials: list[str], output_file: str,chat_message_id: int) -> str:
        # Create a PowerPoint presentation
        prs = Presentation()

        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"Introducing {product_name}"
        subtitle.text = "A Presentation on Key Features, Benefits, and Success Stories"

        # Key Features Slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Key Features"
        content.text = "\n".join([f"- {feature}" for feature in key_features])

        # Benefits Slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Benefits"
        content.text = "\n".join([f"- {benefit}" for benefit in benefits])

        # Customer Testimonials Slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "What Customers Are Saying"
        content.text = "\n".join([f'"{testimonial}"' for testimonial in testimonials])

        # Thank You Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Thank You!"
        subtitle.text = f"Learn more about {product_name} today."

        # Save the presentation
        prs.save(output_file)


        return f"Presentation '{output_file}' has been{chat_message_id} successfully created."


from crewai.tools import BaseTool
from typing import Type
from pydantic import BaseModel, Field


class EmailTemplateCreatorInput(BaseModel):
    """Input schema for EmailTemplateCreator."""
    sender_name: str = Field(..., description="Name of the sender.")
    company_name: str = Field(..., description="Company name for branding.")
    lead_generation_offer: str = Field(..., description="Brief description of the lead generation offer or value proposition.")
    follow_up_context: str = Field(..., description="Context for follow-up (e.g., previous interaction, missed response).")
    re_engagement_reason: str = Field(..., description="Reason for re-engagement (e.g., new features, special offer).")


class EmailTemplateCreatorTool(BaseTool):
    name: str = "Email Template Creator"
    description: str = (
        "Generates three customizable cold email templates for sales scenarios: lead generation, follow-ups, and re-engagement. "
        "Each email includes a subject line, a compelling call to action, and placeholders for personalization."
    )
    args_schema: Type[BaseModel] = EmailTemplateCreatorInput

    def _run(self, sender_name: str, company_name: str, lead_generation_offer: str, follow_up_context: str, re_engagement_reason: str) -> str:
        email_templates = []

        # Lead Generation Email Template
        lead_email = f"""
        **Subject:** Unlock Exclusive Benefits with {company_name}  
        
        Hi {{first_name}},  

        I hope you're doing well! I wanted to introduce you to {company_name} and share how we can help {lead_generation_offer}.  

        If you're interested, I'd love to set up a quick call to discuss how this can benefit you. Let me know a time that works for you!  

        Looking forward to your thoughts.  

        Best,  
        {sender_name}  
        """

        # Follow-Up Email Template
        follow_up_email = f"""
        **Subject:** Following Up on Our Last Conversation  

        Hi {{first_name}},  

        I wanted to follow up regarding {follow_up_context}. I know things get busy, so I wanted to check if you had any questions or if now is a better time to connect.  

        Let me know if you're available for a quick call this week. Looking forward to hearing your thoughts!  

        Best,  
        {sender_name}  
        """

        # Re-Engagement Email Template
        re_engagement_email = f"""
        **Subject:** Exciting Updates from {company_name}!  

        Hi {{first_name}},  

        It’s been a while since we last connected, and I wanted to reach out because {re_engagement_reason}.  

        We’d love to hear your thoughts and explore how we can add value to your business. Let’s reconnect—how about a quick call next week?  

        Looking forward to catching up!  

        Best,  
        {sender_name}  
        """

        email_templates.extend([lead_email, follow_up_email, re_engagement_email])

        return "Email templates successfully created:\n\n" + "\n\n".join(email_templates)
